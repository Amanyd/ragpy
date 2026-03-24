import logging
from collections import defaultdict
from pathlib import Path

from llama_index.core import Document

logger = logging.getLogger(__name__)

_TEXT_RATIO_THRESHOLD = 0.30

_vision_available: bool | None = None

def _check_vision_available() -> bool:
    global _vision_available
    if _vision_available is not None:
        return _vision_available

    import httpx
    from app.config.settings import settings

    try:
        resp = httpx.get(
            f"{settings.ollama_base_url}/api/tags",
            timeout=5.0,
        )
        resp.raise_for_status()
        models = resp.json().get("models", [])
        base_name = settings.ollama_model.split(":")[0]
        for model in models:
            if model.get("name", "").startswith(base_name):
                families = model.get("details", {}).get("families", [])
                _vision_available = "clip" in families or "vision" in families
                logger.info(
                    "vision_probe model=%s available=%s families=%s",
                    settings.ollama_model,
                    _vision_available,
                    families,
                )
                return _vision_available
        logger.warning(
            "vision_probe model=%s not found in Ollama",
            settings.ollama_model,
        )
    except Exception:
        logger.warning("vision_probe_failed — defaulting to no vision")

    _vision_available = False
    return False

def _table_to_markdown(raw_rows: list[list[str]]) -> str:
    if not raw_rows:
        return ""

    def dedup_row(row: list[str]) -> list[str]:
        if not row:
            return row
        out = [row[0]]
        for cell in row[1:]:
            if cell != out[-1]:
                out.append(cell)
        return out

    deduped = [dedup_row(r) for r in raw_rows]
    col_count = max(len(r) for r in deduped)

    def pad(row: list[str]) -> list[str]:
        return (row + [""] * col_count)[:col_count]

    header = pad(deduped[0])
    lines = [
        "| " + " | ".join(header) + " |",
        "|" + "|".join(["---"] * col_count) + "|",
    ]
    for row in deduped[1:]:
        lines.append("| " + " | ".join(
            c.strip().replace("\n", " ") for c in pad(row)
        ) + " |")
    return "\n".join(lines)

def parse_docx(file_path: Path) -> list[Document]:
    from docx import Document as DocxDocument
    from docx.table import Table

    doc = DocxDocument(str(file_path))
    parts: list[str] = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1]

        if tag == "p":
            para = next(
                (p for p in doc.paragraphs if p._element is element), None
            )
            if para and para.text.strip():
                parts.append(para.text.strip())

        elif tag == "tbl":
            table = Table(element, doc)
            raw_rows: list[list[str]] = []
            for row in table.rows:
                raw_rows.append([
                    cell.text.strip().replace("\n", " ")
                    for cell in row.cells
                ])

            if raw_rows:
                parts.append(_table_to_markdown(raw_rows))

    if not parts:
        logger.warning("docx_empty file=%s", file_path.name)
        return []

    text = "\n\n".join(parts)
    logger.info("docx_parsed file=%s len=%d", file_path.name, len(text))
    return [Document(text=text)]

def _extract_slide_text(slide) -> str:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts: list[str] = []

    def process_shape(shape) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                process_shape(s)
            return
        if shape.has_table:
            raw_rows = [
                [
                    cell.text_frame.text.strip().replace("\n", " ")
                    for cell in row.cells
                ]
                for row in shape.table.rows
            ]
            if raw_rows:
                parts.append(_table_to_markdown(raw_rows))
            return
        if shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    indent = "  " * para.level if para.level else ""
                    lines.append(f"{indent}{line}")
            if lines:
                parts.append("\n".join(lines))

    for shape in slide.shapes:
        process_shape(shape)

    return "\n\n".join(parts)

def _extract_slide_notes(slide) -> str:
    try:
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf:
                return tf.text.strip()
    except Exception:
        pass
    return ""

def _rasterize_slide(pptx_path: Path, slide_index: int) -> bytes:
    import io
    import shutil
    import subprocess
    import tempfile

    from pdf2image import convert_from_path

    tmp = Path(tempfile.mkdtemp(prefix="rag_pptx_"))
    try:
        subprocess.run(
            [
                "libreoffice", "--headless",
                "--convert-to", "pdf",
                "--outdir", str(tmp),
                str(pptx_path),
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )
        pdf_path = tmp / (pptx_path.stem + ".pdf")
        if not pdf_path.exists():
            logger.error(
                "libreoffice_pdf_missing file=%s slide=%d",
                pptx_path.name, slide_index,
            )
            return b""

        images = convert_from_path(
            pdf_path,
            dpi=150,
            first_page=slide_index,
            last_page=slide_index,
        )
        if not images:
            return b""

        buf = io.BytesIO()
        images[0].save(buf, format="PNG")
        return buf.getvalue()

    except subprocess.CalledProcessError as e:
        logger.error(
            "libreoffice_failed file=%s slide=%d stderr=%s",
            pptx_path.name, slide_index,
            e.stderr.decode(errors="replace")[:200],
        )
        return b""
    except Exception:
        logger.exception(
            "rasterize_failed file=%s slide=%d", pptx_path.name, slide_index
        )
        return b""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

def _describe_slide_with_vision(png_bytes: bytes, slide_num: int) -> str:
    import base64

    import httpx
    from app.config.settings import settings

    b64 = base64.b64encode(png_bytes).decode()
    payload = {
        "model": settings.ollama_model,
        "prompt": (
            "This is a slide from an educational presentation. "
            "Extract ALL text visible on this slide exactly as written. "
            "Preserve bullet points, numbering, and any table structure. "
            "If there are diagrams or charts, describe what they show."
        ),
        "images": [b64],
        "stream": False,
    }
    try:
        resp = httpx.post(
            f"{settings.ollama_base_url}/api/generate",
            json=payload,
            timeout=180.0,
        )
        resp.raise_for_status()
        return resp.json().get("response", "").strip()
    except Exception:
        logger.exception(
            "vision_inference_failed slide=%d model=%s",
            slide_num, settings.ollama_model,
        )
        return ""

def parse_pptx(file_path: Path) -> list[Document]:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    total_slides = len(prs.slides)

    if total_slides == 0:
        logger.warning("pptx_empty file=%s", file_path.name)
        return []

    slide_texts = [_extract_slide_text(slide) for slide in prs.slides]
    slides_with_text = sum(1 for t in slide_texts if t.strip())
    text_ratio = slides_with_text / total_slides

    vision_supported = _check_vision_available()
    use_vision = text_ratio < _TEXT_RATIO_THRESHOLD and vision_supported

    logger.info(
        "pptx_probe file=%s total=%d with_text=%d ratio=%.2f "
        "vision_supported=%s use_vision=%s",
        file_path.name, total_slides, slides_with_text,
        text_ratio, vision_supported, use_vision,
    )

    if text_ratio < _TEXT_RATIO_THRESHOLD and not vision_supported:
        logger.warning(
            "pptx_image_heavy file=%s ratio=%.2f",
            file_path.name, text_ratio,
        )

    documents: list[Document] = []

    for i, slide in enumerate(prs.slides, 1):
        notes = _extract_slide_notes(slide)

        if use_vision:
            png = _rasterize_slide(file_path, i)
            if not png:
                logger.warning(
                    "pptx_rasterize_empty file=%s slide=%d",
                    file_path.name, i,
                )
                continue
            content = _describe_slide_with_vision(png, i)
        else:
            content = slide_texts[i - 1]
            if not content.strip():
                logger.warning(
                    "pptx_slide_no_text file=%s slide=%d",
                    file_path.name, i,
                )
                continue

        full_text = f"Slide {i}:\n\n{content}"
        if notes:
            full_text += f"\n\nSpeaker notes:\n{notes}"

        documents.append(Document(
            text=full_text,
            metadata={"slide_number": i},
        ))

    logger.info(
        "pptx_parsed file=%s slides_indexed=%d/%d",
        file_path.name, len(documents), total_slides,
    )
    return documents

def parse_pdf(file_path: Path) -> list[Document]:
    from llama_index.readers.file import PDFReader

    documents = PDFReader().load_data(file=file_path)
    logger.info(
        "pdf_parsed file=%s pages=%d", file_path.name, len(documents)
    )
    return documents

_PARSERS = {
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".pdf":  parse_pdf,
}

_SUPPORTED_EXTENSIONS = set(_PARSERS.keys())

def parse_file(file_path: Path) -> list[Document]:
    ext = file_path.suffix.lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {ext!r}")
    return parser(file_path)
